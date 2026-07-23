"""Second deck migration: cover + team slides, native method schematic, chart legibility.

Takes the 9-slide deck produced by restructure_deck_5plus4.py to the 11-slide layout
documented in final-storyboard-5-plus-4.md:

    1 cover · 2 team · 3-7 spoken · 8-11 backup

Also fixes the legibility problem: charts were placed 317-400pt wide, which rendered
their internal type at 3.6-4.4pt. Every chart is now placed at 470pt (the practical
maximum that still leaves a text column) against the regenerated 720pt canvas.

Not idempotent. Run once, against a deck that is CLOSED in PowerPoint.
"""

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent
DECK = ROOT / "the-gammas-w3d5-jaime-draft.pptx"
CHARTS = ROOT / "visuals" / "charts"

INK = RGBColor(0x20, 0x27, 0x29)
MUTED = RGBColor(0x70, 0x5C, 0x40)
BOX_FILL = RGBColor(0xF5, 0xF1, 0xEA)
BOX_LINE = RGBColor(0xC9, 0xBF, 0xB0)

MEMBERS = [
    "Arefeh Lali Dehaghi",
    "Goutham Arcod",
    "Jaime Pineda",
    "Kerem Akyurt",
    "Valeria Moraga",
]
AUTHOR_LINE = " · ".join(MEMBERS)
POD_LINE = "Pod 884 “Ifrit Ras el Hanout” · Megapod Lotus · TA Andrea Buccellato · Project TA Azman Akhter"

# Seven schematic boxes. Numbers copied from visuals/manifest.json and datasets.py:37,
# never retyped from memory -- this diagram is the one place with no automated test.
STEPS = [
    ("1 · Cohort", "336 participants · HCP N-back · 2 WM runs · 360 Glasser ROIs · TR 0.72 s"),
    ("2 · Condition frames", "4 s HRF shift · 312 frames per load · no 0-back/2-back overlap"),
    ("3 · FC per condition", "360 × 360 Pearson correlation, for 0-back and for 2-back"),
    ("4 · Network fingerprint", "12 within + 66 between = 78 values (Cole-Anticevic)"),
    ("5 · Reconfiguration", "fingerprint(2-back) − fingerprint(0-back) → 336 × 78"),
    ("6 · Model", "StandardScaler + RidgeCV, fitted inside each training fold"),
    ("7 · Held-out evaluation",
     "repeated 5-fold CV · fixed holdout · 1000-permutation null · "
     "B→A transfer (301 → 100, identity-disjoint)"),
]

CHART_W, CHART_H = 470.0, 264.4   # 16:9 at the widest that still leaves a text column
CHART_L, CHART_T = 32.4, 98.0
TEXT_L, TEXT_W = 516.0, 171.6


# ---------------------------------------------------------------- helpers
def sid(slide, shape_id):
    return next(s for s in slide.shapes if s.shape_id == shape_id)


def find(slide, prefix):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.startswith(prefix):
            return sh
    raise LookupError(f"no shape starting with {prefix!r}")


def set_lines(shape, lines):
    tf = shape.text_frame
    template = deepcopy(tf.paragraphs[0]._p)
    body = tf._txBody
    for p in list(tf.paragraphs):
        body.remove(p._p)
    for _ in lines:
        body.append(deepcopy(template))
    for para, line in zip(tf.paragraphs, lines):
        runs = para.runs
        runs[0].text = line
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)


def set_text(shape, text):
    set_lines(shape, [text])


def drop(shape):
    shape._element.getparent().remove(shape._element)


def place(shape, left=None, top=None, width=None, height=None):
    for attr, val in (("left", left), ("top", top), ("width", width), ("height", height)):
        if val is not None:
            setattr(shape, attr, Pt(val))


def unique_ids(slide):
    seen, nxt = set(), max((sh.shape_id for sh in slide.shapes), default=1) + 1
    for sh in slide.shapes:
        if sh.shape_id in seen:
            sh._element._nvXxPr.cNvPr.set("id", str(nxt))
            nxt += 1
        seen.add(sh.shape_id)


def textbox(slide, left, top, width, height, text, size, bold=False,
            color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Pt(left), Pt(top), Pt(width), Pt(height))
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Poppins"
    run.font.color.rgb = color
    return box


def chrome(slide, kicker, page, total=11):
    """Brand footer, page marker and status kicker, copied from an existing slide."""
    set_text(sid(slide, 4), kicker)
    set_text(sid(slide, 3), f"{page} / {total}")


def reset_chart(slide, slug):
    """Drop every embedded picture and re-insert the CURRENT chart file.

    The deck stores its own copy of each PNG, so repositioning alone would keep the
    pre-regeneration image. Re-inserting from disk is what actually picks up the
    720pt-canvas rendering.
    """
    for sh in [s for s in slide.shapes if s.shape_type == 13]:
        drop(sh)
    source = CHARTS / f"{slug}.png"
    if not source.is_file():
        raise FileNotFoundError(source)
    slide.shapes.add_picture(str(source), Pt(CHART_L), Pt(CHART_T),
                             Pt(CHART_W), Pt(CHART_H))


# ---------------------------------------------------------------- pass A
def pass_a():
    """Add cover + team as new slides, then move them to the front. Add-only."""
    prs = Presentation(str(DECK))
    assert len(prs.slides._sldIdLst) == 9, "expected the 9-slide deck"

    donor = prs.slides[0]            # intro: has brand footer, page marker, kicker, title
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    for _ in range(2):
        new = prs.slides.add_slide(blank)
        for ph in list(new.shapes):
            drop(ph)
        for shape_id in (2, 3, 4, 5):   # brand, page, kicker, title
            new.shapes._spTree.append(deepcopy(sid(donor, shape_id)._element))

    cover, team = prs.slides[9], prs.slides[10]

    # ---- cover
    chrome(cover, "THE GAMMAS · NMA COMPUTATIONAL NEUROSCIENCE 2026", 1)
    title = sid(cover, 5)
    place(title, left=40.2, top=112.0, width=639.6, height=104.0)
    set_lines(title, ["Functional Connectivity Reconfiguration",
                      "in N-back Working Memory"])
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
    textbox(cover, 40.2, 232.0, 639.6, 20, AUTHOR_LINE, 12, color=INK)
    textbox(cover, 40.2, 254.0, 639.6, 18,
            "Ifrit Ras el Hanout — The Gammas · NMA Computational Neuroscience 2026", 10,
            color=MUTED)
    textbox(cover, 40.2, 274.0, 639.6, 18,
            "HCP N-back · participant-level prediction", 10, color=MUTED)
    cover.notes_slide.notes_text_frame.text = (
        "Portada. No consume minuto; se muestra mientras se saluda y se presenta al equipo."
    )

    # ---- team
    chrome(team, "THE GAMMAS · NMA COMPUTATIONAL NEUROSCIENCE 2026", 2)
    set_text(sid(team, 5), "Meet our team")
    # 3-over-2 monogram grid; names only -- template roles are dummy text, photos are stock
    row1_x, row2_x = (60.0, 285.0, 510.0), (172.0, 397.0)
    for i, member in enumerate(MEMBERS):
        x = row1_x[i] if i < 3 else row2_x[i - 3]
        y = 130.0 if i < 3 else 248.0
        circle = team.shapes.add_shape(MSO_SHAPE.OVAL, Pt(x + 37), Pt(y), Pt(76), Pt(76))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BOX_FILL
        circle.line.color.rgb = BOX_LINE
        circle.line.width = Pt(1)
        tf = circle.text_frame
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = "".join(p[0] for p in member.split()[:2]).upper()
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.name = "Poppins"
        run.font.color.rgb = MUTED
        textbox(team, x, y + 82, 150, 18, member, 11, bold=True, align=PP_ALIGN.CENTER)
    textbox(team, 32.4, 352.0, 655.2, 18, POD_LINE, 9, color=MUTED, align=PP_ALIGN.CENTER)
    team.notes_slide.notes_text_frame.text = (
        "Solo nombres. Los roles de la plantilla (Actor/Doctor/Accountant/Chef/Journalist) son "
        "texto de relleno y se han eliminado; los retratos de la plantilla son de stock y no "
        "somos nosotros. Si el equipo quiere rol real por persona, se añade con su confirmación."
    )

    for slide in (cover, team):
        unique_ids(slide)

    lst = prs.slides._sldIdLst
    entries = list(lst)
    for e in entries:
        lst.remove(e)
    for i in [9, 10, 0, 1, 2, 3, 4, 5, 6, 7, 8]:
        lst.append(entries[i])

    prs.save(str(DECK))
    print("pass A: cover + team added and moved to front")


# ---------------------------------------------------------------- pass B
def pass_b():
    """Shape-level edits: schematic, inset removal, chart enlargement, paging."""
    prs = Presentation(str(DECK))
    assert len(prs.slides._sldIdLst) == 11, "expected 11 slides after pass A"
    (_cover, _team, intro, method, primary,
     turn, conclusion, bk_dir, bk_val, bk_future, bk_refs) = prs.slides

    # ---- intro: roster moved to the team slide
    chrome(intro, "SPOKEN · 1 / 5", 3)
    for prefix in ("Arefeh Lali Dehaghi ·", "HCP N-back · participant-level"):
        try:
            drop(find(intro, prefix))
        except LookupError:
            pass

    # ---- method: replace the feature-construction chart with a native schematic
    chrome(method, "SPOKEN · 2 / 5", 4)
    set_text(sid(method, 5), "One pipeline, from scans to a held-out prediction")
    for sh in list(method.shapes):
        if sh.shape_type == 13 or (sh.has_text_frame and sh.text_frame.text.startswith("•  Cohort B")):
            drop(sh)

    def step_box(left, top, width, height, label, detail):
        box = method.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Pt(left), Pt(top), Pt(width), Pt(height))
        box.fill.solid()
        box.fill.fore_color.rgb = BOX_FILL
        box.line.color.rgb = BOX_LINE
        box.line.width = Pt(1)
        box.adjustments[0] = 0.08
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(6)
        tf.margin_top = tf.margin_bottom = Pt(4)
        head = tf.paragraphs[0]
        run = head.add_run()
        run.text = label
        run.font.size, run.font.bold, run.font.name = Pt(10.5), True, "Poppins"
        run.font.color.rgb = INK
        body = tf.add_paragraph()
        run = body.add_run()
        run.text = detail
        run.font.size, run.font.name = Pt(8), "Poppins"
        run.font.color.rgb = MUTED

    row1_w, row1_gap = 152.0, 15.7
    for i in range(4):
        step_box(32.4 + i * (row1_w + row1_gap), 104.0, row1_w, 76.0, *STEPS[i])
    for i in range(2):
        step_box(32.4 + i * (311.6 + 32.0), 194.0, 311.6, 68.0, *STEPS[4 + i])
    step_box(32.4, 276.0, 655.2, 62.0, *STEPS[6])
    textbox(method, 32.4, 346.0, 655.2, 16,
            "All splits hold out participants; “±” = split SD, not CI.", 8.5, color=MUTED)

    # ---- primary: drop the illegible inset, re-embed the regenerated transfer chart
    chrome(primary, "SPOKEN · 3 / 5", 5)
    reset_chart(primary, "identity-disjoint-transfer")
    try:
        drop(find(primary, "Primary repeated-CV (compact reference)"))
    except LookupError:
        pass
    for prefix, top in (("B→A transfer", 104.0), ("r = 0.398", 122.0),
                        ("301 B-only", 152.0), ("Primary repeated CV", 200.0),
                        ("r = 0.366", 218.0)):
        place(find(primary, prefix), left=TEXT_L, top=top, width=TEXT_W)

    # ---- turn: re-embed chart, surface the no-winner caveat as real slide text
    chrome(turn, "SPOKEN · 4 / 5", 6)
    reset_chart(turn, "activation-robustness")
    textbox(turn, CHART_L, 366.0, CHART_W, 14, "No biological winner is claimed.", 9,
            bold=True, color=MUTED)

    # ---- conclusion
    chrome(conclusion, "SPOKEN · 5 / 5 · KEEP VISIBLE FOR Q&A", 7)

    # ---- backups
    for slide, page, kicker, slug in (
        (bk_dir, 8, "BACKUP · OPEN ONLY IF ASKED", "segregation-refinement"),
        (bk_val, 9, "BACKUP · OPEN ONLY IF ASKED", "null-and-holdout"),
        (bk_future, 10, "BACKUP · OPEN ONLY IF ASKED", None),
        (bk_refs, 11, "BACKUP · Q&A ONLY", None),
    ):
        chrome(slide, kicker, page)
        if slug:
            reset_chart(slide, slug)

    for slide in prs.slides:
        unique_ids(slide)

    prs.save(str(DECK))
    print("pass B: schematic, inset removed, charts enlarged, paging = /11")


if __name__ == "__main__":
    lock = DECK.parent / f"~${DECK.name}"
    if lock.exists():
        raise SystemExit(f"PowerPoint has the deck open ({lock.name}). Close it and re-run.")
    pass_a()
    pass_b()
