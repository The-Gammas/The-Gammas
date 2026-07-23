"""One-shot restructure of the W3D5 deck to 5 spoken + 4 backup slides.

Aligns slide architecture with the 2025 NMA example decks (fused title/team opener,
dedicated surprise slot, one dominant figure per spoken slide, backup tail) without
changing any scientific content. Every number kept here already exists in
pipeline/02 and is recorded in visuals/manifest.json.

Slide count stays at 9 and no slide part is created or destroyed: the silent cover is
REPURPOSED into the new validation-backup slide. Deleting + adding slides in one
python-pptx session reassigns a colliding partname and silently overwrites a slide.
"""

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(
    "/Users/jaimepm/Library/Mobile Documents/com~apple~CloudDocs/Wiki/"
    "Life Long Learning/NeuroAcademy/project/fmri/the-gammas/manuscript/slides"
)
DECK = ROOT / "the-gammas-w3d5-jaime-draft.pptx"
CHARTS = ROOT / "visuals" / "charts"

prs = Presentation(str(DECK))

COVER, INTRO, METHOD, PRIMARY, SEGREG, TURN, CONCL, FUTURE, REFS = range(9)


# ---------------------------------------------------------------- helpers
def sid(slide, shape_id):
    return next(s for s in slide.shapes if s.shape_id == shape_id)


def set_lines(shape, lines):
    """Replace paragraph text, cloning paragraph 0's run formatting for each line."""
    tf = shape.text_frame
    template = deepcopy(tf.paragraphs[0]._p)
    body = tf._txBody
    for p in list(tf.paragraphs):
        body.remove(p._p)
    for _ in lines:
        body.append(deepcopy(template))
    for para, line in zip(tf.paragraphs, lines):
        runs = para.runs
        runs[0].text = line
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)


def set_text(shape, text):
    set_lines(shape, [text])


def drop(shape):
    shape._element.getparent().remove(shape._element)


def place(shape, left=None, top=None, width=None, height=None):
    for attr, val in (("left", left), ("top", top), ("width", width), ("height", height)):
        if val is not None:
            setattr(shape, attr, Pt(val))


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def reorder(order):
    lst = prs.slides._sldIdLst
    entries = list(lst)
    for e in entries:
        lst.remove(e)
    for i in order:
        lst.append(entries[i])


# ------------------------------------------------- 1. intro (fused opener)
cover = prs.slides[COVER]
brand_line = deepcopy(sid(cover, 7)._element)
subtitle_line = deepcopy(sid(cover, 8)._element)

intro = prs.slides[INTRO]
set_text(sid(intro, 4), "SPOKEN · 1 / 5")
set_text(
    sid(intro, 5),
    "Does load-related brain connectivity predict working-memory performance?",
)
set_text(sid(intro, 7), "Pattern hypothesis")
set_text(
    sid(intro, 8),
    "A 78-feature FC fingerprint of the 2-back − 0-back change predicts performance.",
)
drop(sid(intro, 9))
drop(sid(intro, 10))
place(sid(intro, 14), top=190.8)

intro.shapes._spTree.append(brand_line)
intro.shapes._spTree.append(subtitle_line)
roster, subtitle = intro.shapes[-2], intro.shapes[-1]
place(roster, left=32.4, top=258.0, width=655.2, height=20.0)
set_text(
    roster,
    "Arefeh Lali Dehaghi · Goutham Arcod · Jaime Pineda · Kerem Akyurt · Valeria Moraga",
)
place(subtitle, left=32.4, top=278.0, width=655.2, height=20.0)
set_text(subtitle, "HCP N-back · participant-level prediction")

set_notes(
    intro,
    "Abrimos con la pregunta y con las DOS cosas que predijimos. Importa decir las dos al "
    "principio: una se confirma y la otra se matiza, y el relato solo se entiende si la "
    "audiencia oyó ambas antes de ver resultados.\n\n"
    "NO decir que la hipótesis evolucionó, cambió o fue reemplazada. Pre-especificamos dos "
    "hipótesis; la evidencia confirmó una y matizó la otra.\n\n"
    "Caveat: 'reconfiguration' es una diferencia entre dos FC agregadas por condición, no "
    "conectividad dinámica. La hipótesis direccional es la más restrictiva: exige un signo y "
    "comprime una estructura multivariada en un escalar.",
)

# ------------------------------------------------- 2. method
method = prs.slides[METHOD]
set_text(sid(method, 4), "SPOKEN · 2 / 5")
set_notes(
    method,
    "Scaler y RidgeCV se ajustan dentro de cada fold; nunca se separan runs o frames de la "
    "misma persona. B estima el modelo; A permite comprobar transferencia sin identidades "
    "compartidas.\n\n"
    "Evidencia: pipeline/02 §3–§4 (B=339 task / 336 analíticos, A=100, 360 ROIs, 12 redes, "
    "78 features, 4 s HRF, 312 frames). Chart: manifest.json → feature-construction (celdas 6, 8).\n"
    "Caveat: la FC de tarea basada en Pearson puede reflejar coactivación evocada; no demuestra "
    "comunicación interregional ni causalidad.",
)

# ------------------------------------------------- 3. primary result (trimmed)
primary = prs.slides[PRIMARY]
set_text(sid(primary, 4), "SPOKEN · 3 / 5")
set_text(
    sid(primary, 5),
    "The FC pattern predicted performance—and transferred to a separate cohort",
)
for shape_id in (11, 12, 13, 14):  # fixed holdout + full-refit null -> backup
    drop(sid(primary, shape_id))
place(sid(primary, 15), top=100.8)
place(sid(primary, 16), top=119.5)
place(sid(primary, 17), top=149.8)
place(sid(primary, 9), top=198.7)
place(sid(primary, 10), top=217.4)
set_notes(
    primary,
    "El efecto primario es la media de CV repetida, .366 ± .024. La transferencia usa modelos "
    "entrenados solo en B y aplicados a A sin identidades compartidas: no es validación en otro "
    "sitio. Si preguntan por el holdout fijo o por el null, están en la slide de backup.\n\n"
    "Transición: establecida la predicción, toca revisar lo que NO salió como esperábamos.\n\n"
    "Evidencia: pipeline/02 §5–§6, paneles 9B/9C. Charts: manifest.json → "
    "identity-disjoint-transfer (celda 14), primary-repeated-cv (celda 10).\n"
    "Caveat: la SD resume sensibilidad entre 20 particiones solapadas, no incertidumbre "
    "poblacional. Cohortes HCP hermanas no son generalización entre sitios; el parentesco no "
    "está modelado.",
)

# ------------------------------------------------- 4. the turn (activation)
turn = prs.slides[TURN]
set_text(sid(turn, 4), "SPOKEN · 4 / 5")
set_text(sid(turn, 5), "The prediction held—but a simpler regional signal predicted better")
set_lines(
    sid(turn, 19),
    [
        "Post hoc, unmatched comparison: 360 regional activation vs 78 network FC features. "
        "Activation = mean BOLD(2-back) − mean BOLD(0-back), not a GLM beta. "
        "Vascular reactivity (CVR) not controlled."
    ],
)
turn.shapes._spTree.append(deepcopy(sid(turn, 7)._element))
bridge = turn.shapes[-1]
place(bridge, left=32.4, top=328.0, width=399.6, height=46.0)
set_lines(
    bridge,
    [
        "Direction, as predicted—but only at group level:",
        "segregation fell under load (Δ = −0.0236; p = 3.45 × 10⁻⁵), yet larger shifts "
        "did not predict better performance (r = −0.105; p = .054).",
    ],
)
for para in bridge.text_frame.paragraphs:
    for run in para.runs:
        run.font.size = Pt(9)
bridge.text_frame.paragraphs[0].runs[0].font.bold = True

set_notes(
    turn,
    "Este es el giro, y hay que contarlo en dos tiempos.\n\n"
    "UNO: la dirección que predijimos existe — la segregación baja con la carga — pero como "
    "predictor individual es débil, así que la hipótesis direccional queda MATIZADA, no "
    "confirmada. La figura completa está en backup.\n\n"
    "DOS: al probar si la señal era específica de conectividad, un contraste de activación "
    "regional mucho más simple predijo mejor. Es una prueba POST HOC de robustez, no una "
    "competición biológica justa: 360 rasgos de activación frente a 78 de FC, sin igualar "
    "dimensionalidad. Los dos ΔR² usan folds emparejados y la regla de 2 SD es una heurística "
    "de sensibilidad al split, no un test formal de superioridad.\n\n"
    "La activación NO borra el resultado FC; impide afirmar que el mecanismo sea específico de "
    "conectividad. Nuestros controles (DVARS, acc_0bk) no capturan reactividad cerebrovascular "
    "(CVR, señalada por Goutham el 22 jul; Logothetis 2008), así que ese hueco queda declarado.\n\n"
    "Evidencia: pipeline/02 §7 + bloque de controles + §8. Chart: manifest.json → "
    "activation-robustness (celdas 18, 20); cifras de segregación de segregation-refinement "
    "(celda 22).\n"
    "Caveat: el centrado por run hace 0-back, 2-back y su contraste fuertemente colineales — un "
    "solo eje de activación visto de tres formas, no un rasgo independiente de carga.",
)

# ------------------------------------------------- 5. conclusion
conclusion = prs.slides[CONCL]
set_text(sid(conclusion, 4), "SPOKEN · 5 / 5 · KEEP VISIBLE FOR Q&A")
set_lines(
    sid(conclusion, 8),
    [
        "•  Pattern hypothesis: a 78-feature FC difference predicts unseen 2-back accuracy.",
        "•  The model transfers across identity-disjoint same-HCP cohorts.",
    ],
)
set_lines(
    sid(conclusion, 11),
    [
        "•  Directional hypothesis: segregation fell under load, but larger shifts did not "
        "predict better performance.",
        "•  Reconfiguration showed no clear gain beyond 0-back FC.",
        "•  FC added no clear gain over activation under the current unmatched comparison.",
    ],
)
set_lines(
    sid(conclusion, 14),
    [
        "•  Is predictive information connectivity-specific rather than shared with task "
        "activation?",
        "•  Vascular reactivity (CVR) is not controlled in the activation benchmark.",
    ],
)
set_notes(
    conclusion,
    "Cerrar NOMBRANDO las dos predicciones de la Slide 1: la de patrón SE SOSTIENE, la "
    "direccional SE MATIZA. Frase final: 'predictive signal survives; connectivity-specific "
    "mechanism remains unresolved.'\n\n"
    "Mantener esta slide proyectada durante las preguntas; las de backup solo se abren si hacen "
    "falta.\n\n"
    "Evidencia: pipeline/02 §11–§12 (tabla claim → result → figure → limitation).\n"
    "Caveat: estudio observacional de una sola tarea: no identifica causalidad, conectividad "
    "dinámica, beneficio adaptativo ni especificidad biológica.",
)

# ------------------------------------------------- 6. backup: directional
backup_dir = prs.slides[SEGREG]
set_text(sid(backup_dir, 4), "BACKUP · OPEN ONLY IF ASKED")
set_text(sid(backup_dir, 5), "Group direction was real; the individual link was weak")
set_notes(
    backup_dir,
    "Bajada del deck hablado para que la slide del giro conserve una sola figura dominante. "
    "Abrir si preguntan qué significa exactamente que 'la hipótesis direccional quedó matizada'.\n\n"
    "Evidencia: pipeline/02 §8; manifest.json → segregation-refinement (celda 22).\n"
    "Caveat: la segregación de sistema es un escalar tipo Chan, no modularidad de Newman, y no "
    "sustituye al patrón multivariado de 78 rasgos. No reutilizar la magnitud −0.048 del abstract.",
)

# ---------------------------- 7. cover REPURPOSED as validation backup slide
validation = prs.slides[COVER]
for shape in list(validation.shapes):
    if shape.shape_id not in (2, 3):  # keep brand footer + page marker
        drop(shape)
for shape_id in (4, 5, 7, 8, 9):  # kicker, title, bullets, rule, caveat
    validation.shapes._spTree.append(deepcopy(sid(backup_dir, shape_id)._element))

set_text(sid(validation, 4), "BACKUP · OPEN ONLY IF ASKED")
set_text(sid(validation, 5), "The checks behind the primary result")
set_lines(
    sid(validation, 7),
    [
        "•  Fixed holdout: r = 0.312 in 67 unseen participants",
        "•  Full-refit null (seed 42): r = 0.405; p = 1/1001 ≈ .001",
        "•  B→A A-label permutation, fixed B predictions: p = 1/1001 ≈ .001",
        "•  Reconfiguration over 0-back FC: ΔR² = +0.0344 ± 0.0225 → no clear gain",
    ],
)
set_text(
    sid(validation, 9),
    "The permutation p belongs only to seed-42 r = 0.405; the holdout is a separate split.",
)
validation.shapes.add_picture(
    str(CHARTS / "null-and-holdout.png"),
    Pt(32.4), Pt(108.0), Pt(388.8), Pt(218.7),
)
set_notes(
    validation,
    "Bajada de la slide de resultado primario para que lleve una figura dominante. Abrir ante "
    "'¿cómo sabéis que no es azar?' o '¿y el holdout?'.\n\n"
    "Evidencia: pipeline/02 §5–§7; manifest.json → null-and-holdout (celda 12), "
    "incremental-fc-test (celda 18).\n"
    "Caveat: la regla de decisión de 2 SD es una heurística de sensibilidad al split, no un test "
    "formal de superioridad ni de equivalencia.",
)

# ------------------------------------------------- 8. remaining backup kickers
set_text(sid(prs.slides[FUTURE], 4), "BACKUP · OPEN ONLY IF ASKED")
set_text(sid(prs.slides[REFS], 4), "BACKUP · Q&A ONLY")

# ------------------------------------------------- 9. final order + paging
reorder([INTRO, METHOD, PRIMARY, TURN, CONCL, SEGREG, COVER, FUTURE, REFS])

for i, slide in enumerate(prs.slides):
    set_text(sid(slide, 3), f"{i + 1} / 9")

prs.save(str(DECK))
print("saved")
